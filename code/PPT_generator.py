from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.dml.color import RGBColor
import json
import os
from utils.os_helpers import resize_image
from random_generator import modify_style
import shutil
import ast
import random


def add_style_pertubations(textbox, text, style_indexes, default_style, special_style):
    tf = textbox.text_frame
    words = text.split()
    tf.auto_size = True
    tf.word_wrap = True
    current_index = 0
    current_run = None
    para = tf.add_paragraph()
    for i, word in enumerate(words):
        word += ' '
        for index in style_indexes:
            start = index[0]
            end = index[-1]
            if i >= start and i <= end:
                current_run = None
                if current_run != None:
                    current_run = para.add_run()
                    current_run.text = word
                else:
                    current_run = para.add_run()
                    current_run.text = word
                    current_run.font.bold = special_style.get('bold', False)
                    current_run.font.color.rgb = RGBColor(special_style['font_color']['r'], special_style['font_color']['g'], special_style['font_color']['b'])                    
                    current_run.font.italic = special_style.get('italics', False)
                    current_run.font.underline = special_style.get('underlined', False)
                    current_run.font.size = Pt(special_style.get('font_size', 16))
                    current_run.font.name = special_style.get('font_name', 'Arial')
                break
        else:
            current_run = para.add_run()
            current_run.text = word
            current_run.font.color.rgb = RGBColor(default_style['font_color']['r'], default_style['font_color']['g'], default_style['font_color']['b'])                    
            current_run.font.bold = default_style.get('bold', False)
            current_run.font.italic = default_style.get('italics', False)
            current_run.font.underline = default_style.get('underlined', False)
            current_run.font.size = Pt(default_style.get('font_size', 16))
            current_run.font.name = default_style.get('font_name', 'Arial')
                
        current_index += len(word) + 1  
        
    current_index += len(word) + 1  # Account for spaces between words

# Ensure the last part of the text is added
    if current_index < len(text):
        current_run.text += text[current_index:]
        current_run = None
# Reset the font properties for the last run
    current_run.font.bold = default_style.get('bold', False)
    current_run.font.italic = default_style.get('italic', False)
    current_run.font.underline = default_style.get('underlined', False)
    current_run.font.size = Pt(default_style.get('font_size', 16))
    current_run.font.name = default_style.get('font_name', 'Arial')


class Element:
    def __init__(self, content, style, bounding_box):
        self.content = content
        self.style = style
        self.bounding_box = bounding_box
    
    def apply_font_style_on_run(self, run):
        if 'font_name' in self.style:
            run.font.name = self.style['font_name']
        if 'font_size' in self.style:
            run.font.size = Pt(self.style['font_size'])
        if 'font_color' in self.style:
            run.font.color.rgb = RGBColor(self.style['font_color']['r'], self.style['font_color']['g'], self.style['font_color']['b'])
        if 'bold' in self.style:
            run.font.bold = self.style['bold']
        if 'italics' in self.style:
            run.font.italic = self.style['italics']
        if 'underlined' in self.style:
            run.font.underline = self.style['underlined']
        
        if 'h_align' in self.style:
            if self.style['h_align'] == 'center':
                run.alignment = PP_ALIGN.CENTER
            elif self.style['h_align'] == 'left':
                run.alignment = PP_ALIGN.LEFT
            elif self.style['h_align'] == 'right':
                run.alignment = PP_ALIGN.RIGHT
            else:
                run.alignment = PP_ALIGN.JUSTIFY
        
        if 'v_align' in self.style:
            if self.style['v_align'] == 'top':
                run.vertical_anchor = MSO_ANCHOR.TOP
            elif self.style['v_align'] == 'middle':
                run.vertical_anchor = MSO_ANCHOR.MIDDLE
            else:
                run.vertical_anchor = MSO_ANCHOR.BOTTOM

    
    def apply_font_style(self, shape):
        if 'font_name' in self.style:
            shape.text_frame.paragraphs[0].font.name = self.style['font_name']
        if 'font_size' in self.style:
            shape.text_frame.paragraphs[0].font.size = Pt(self.style['font_size'])
        if 'font_color' in self.style:
            shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(self.style['font_color']['r'], self.style['font_color']['g'], self.style['font_color']['b'])
        if 'bold' in self.style:
            shape.text_frame.paragraphs[0].font.bold = self.style['bold']
        else:
            shape.text_frame.paragraphs[0].font.bold = False 
        if 'italics' in self.style:
            shape.text_frame.paragraphs[0].font.italic = self.style['italics']
        else:
            shape.text_frame.paragraphs[0].font.italic = False 
        if 'underlined' in self.style:
            shape.text_frame.paragraphs[0].font.underline = self.style['underlined']
        else:
            shape.text_frame.paragraphs[0].font.underline = False
        if 'h_align' in self.style:
            if self.style['h_align'] == 'center':
                shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            elif self.style['h_align'] == 'left':
                shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            elif self.style['h_align'] == 'right':
                shape.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
            else:
                shape.text_frame.paragraphs[0].alignment = PP_ALIGN.JUSTIFY
        
        if 'v_align' in self.style:
            if self.style['v_align'] == 'top':
                shape.text_frame.paragraphs[0].vertical_anchor = MSO_ANCHOR.TOP
            elif self.style['v_align'] == 'middle':
                shape.text_frame.paragraphs[0].vertical_anchor = MSO_ANCHOR.MIDDLE
            else:
                shape.text_frame.paragraphs[0].vertical_anchor = MSO_ANCHOR.BOTTOM

    
    def position_element(self, shape):
        shape.left = Inches(self.bounding_box[0])
        shape.top = Inches(self.bounding_box[1])
        shape.width = Inches(self.bounding_box[2])
        shape.height = Inches(self.bounding_box[3])

        if shape.has_text_frame:
            shape.text_frame.margin_left = Pt(10)
            shape.text_frame.margin_right = Pt(10)
            shape.text_frame.margin_top = Pt(10)
            shape.text_frame.margin_bottom = Pt(10)


    def render(self, slide):
        raise NotImplementedError("Subclasses must implement this method")

class Title(Element):
    def __init__(self, content, style, bounding_box):
        super().__init__(content, style, bounding_box)
    def render(self, slide):
        title_shape = slide.shapes.title
        title_shape.text = self.content
        self.apply_font_style(title_shape)
        self.position_element(title_shape)

    def clean_up(self, slide):
        pass

class Description(Element):
    def __init__(self, content, style, bounding_box, phrases):
        super().__init__(content, style, bounding_box)
        self.phrases = phrases
        self.special_style = modify_style(style)

    def render(self, slide):
        left, top, width, height = self.bounding_box
        textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        add_style_pertubations(textbox, self.content, self.phrases, self.style, self.special_style)
        self.textbox = textbox

class Reference(Element):
    def __init__(self, content, style, bounding_box):
        super().__init__(content, style, bounding_box)
    def render(self, slide):
        left, top, width, height = self.bounding_box
        textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        textbox.text = self.content
        self.apply_font_style(textbox)
        textbox.text_frame.auto_size = True
        textbox.text_frame.word_wrap = True 
        self.position_element(textbox)
        self.textbox = textbox

class Enumeration(Description):
    def __init__(self, content, style, bounding_box, heading, phrases):
        super().__init__(content, style, bounding_box, phrases)
        self.heading = heading

    def render(self, slide):
        if self.heading != None:
            left, top, width, height = self.heading['xmin'], self.heading['ymin'], self.heading['width'], self.heading['height']
            enum_heading = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
            enum_heading.text = self.heading['value']
            self.apply_font_style(enum_heading)
            enum_heading.text_frame.paragraphs[0].font.size = Pt(self.heading['style']['font_size'])
            enum_heading.text_frame.paragraphs[0].font.bold = self.heading['style']['bold']
            enum_heading.text_frame.paragraphs[0].font.underline = self.heading['style']['underlined']
            enum_heading.text_frame.auto_size = True
            enum_heading.text_frame.word_wrap = True
            # enum_heading.text_frame.paragraphs[0].alignment = PP_ALIGN.JUSTIFY
        enum_shape = slide.shapes.placeholders[1]
        enum_tf = enum_shape.text_frame
        if self.content != []:            
            enum_tf.text = self.content[0]
            self.apply_font_style(enum_shape)
            # enum_tf.paragraphs[0].paragraph_format.alignment = MSO_ANCHOR.JUSTIFY
            for i, pt_text in enumerate(self.content):
                if i>0:
                    if isinstance(pt_text, str):
                        p = enum_tf.add_paragraph()
                        run = p.add_run()
                        run.text = pt_text
                        self.apply_font_style_on_run(run)
                        if random.random() > 0.5:
                            p.level = 1
                            run.font.size = Pt(self.style['font_size'] - random.randint(0, 2))
                            if self.style['font_color']['r'] == 0:
                                run.font.color.rgb = RGBColor(169, 169, 169)
                            else:
                                run.font.color.rgb = RGBColor(211, 211, 211)

                    # elif isinstance(pt_text, list):
                    #     for sub_pt in pt_text:
                    #         s_p = enum_tf.add_paragraph()
                    #         sub_run = s_p.add_run()
                    #         sub_run.text = sub_pt
                    #         s_p.level = 1
                    #         self.apply_font_style_on_run(sub_run)
                            # s_p.paragraph_format.alignment = MSO_ANCHOR.JUSTIFY
                    else:
                        raise Exception("Invalid Enumeration format")
                
        enum_tf.auto_size = True
        enum_tf.word_wrap = True
        MAX_LINES = 5
        if len(enum_tf.paragraphs) > MAX_LINES:
            for paragraph in enum_tf.paragraphs[MAX_LINES:]:
                paragraph.text = paragraph.text[:paragraph.text.rfind(' ')] + "..."
        self.position_element(enum_shape)
        self.enum_tf = enum_tf
            

class Figure(Element):
    def __init__(self, content, style, bounding_box, caption):
        super().__init__(content, style, bounding_box)
        self.caption = caption
        self.content = content
        
    def render(self, slide):
        left, top, width, height = self.bounding_box
        resized_img_path, n_w, n_h = resize_image(self.content, width, height)
        if self.caption != None:
            left_c, top_c, width_c, height_c = self.caption['xmin'], self.caption['ymin'], self.caption['width'], self.caption['height']
            if top_c > top:
                cap_shape = slide.shapes.add_textbox(Inches(left_c), Inches(top_c - (height - n_h)/2), Inches(width_c), Inches(height_c))
            else:
                cap_shape = slide.shapes.add_textbox(Inches(left_c), Inches(top_c + (height - n_h)/2), Inches(width_c), Inches(height_c))
            cap_shape.text = self.caption['value']
            self.apply_font_style(cap_shape)
            cap_shape.text_frame.auto_size = True
            cap_shape.text_frame.word_wrap = True
            cap_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        img = slide.shapes.add_picture(resized_img_path, Inches(left - (n_w - width)/2), Inches(top - (n_h - height)/2), Inches(n_w), Inches(n_h))
        
        self.image = img 

class Equation(Element):
    def __init__(self, content, style, bounding_box, caption):
        super().__init__(content, style, bounding_box)
        self.caption = caption
        self.content = content

    def render(self, slide):
        left, top, width, height = self.bounding_box
        resized_img_path, n_w, n_h = resize_image(self.content, width, height)
        if self.caption != None:
            left_c, top_c, width_c, height_c = self.caption['xmin'], self.caption['ymin'], self.caption['width'], self.caption['height']
            if top_c > top:
                cap_shape = slide.shapes.add_textbox(Inches(left_c), Inches(top_c - (height - n_h)/2), Inches(width_c), Inches(height_c))
            else:
                cap_shape = slide.shapes.add_textbox(Inches(left_c), Inches(top_c + (height - n_h)/2), Inches(width_c), Inches(height_c))
            cap_shape.text = self.caption['value']
            self.apply_font_style(cap_shape)
            cap_shape.text_frame.auto_size = True
            cap_shape.text_frame.word_wrap = True
            cap_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        img = slide.shapes.add_picture(resized_img_path, Inches(left - (n_w - width)/2), Inches(top - (n_h - height)/2), Inches(n_w), Inches(n_h))
        self.image = img 


class Table(Element):
    def __init__(self, content, style, bounding_box, caption, tbl_cnt):
        super().__init__(content, style, bounding_box)
        self.caption = caption
        self.content = content
        if tbl_cnt:
            self.tbl_cnt = ast.literal_eval(tbl_cnt)

    def render(self, slide):
        left, top, width, height = self.bounding_box
        if not self.content is None:
            resized_img_path, n_w, n_h = resize_image(self.content, width, height)
            if self.caption != None:
                left_c, top_c, width_c, height_c = self.caption['xmin'], self.caption['ymin'], self.caption['width'], self.caption['height']
                if top_c > top:
                    cap_shape = slide.shapes.add_textbox(Inches(left_c), Inches(top_c - (height - n_h)/2), Inches(width_c), Inches(height_c))
                else:
                    cap_shape = slide.shapes.add_textbox(Inches(left_c), Inches(top_c + (height - n_h)/2), Inches(width_c), Inches(height_c))
                cap_shape.text = self.caption['value']
                self.apply_font_style(cap_shape)
                cap_shape.text_frame.auto_size = True
                cap_shape.text_frame.word_wrap = True
                cap_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            img = slide.shapes.add_picture(resized_img_path, Inches(left - (n_w - width)/2), Inches(top - (n_h - height)/2), Inches(n_w), Inches(n_h))
            self.image = img 
        else:
            rows = len(self.tbl_cnt)
            columns = len(self.tbl_cnt[0])
            x, y, cx, cy = Inches(left), Inches(top), Inches(width), Inches(height)
            shape = slide.shapes.add_table(rows, columns, x, y, cx, cy)
            table = shape.table
            for i in range(rows):
                for j in range(columns):
                    cell = table.cell(i, j)
                    cell.text = self.tbl_cnt[i][j]
            self.table = table
            
        
class Footer(Element):
    def __init__(self, content, style, bounding_box, location):
        super().__init__(content, style, bounding_box)
        self.location = location
    def render(self, slide):
        left, top, width, height = self.bounding_box
        textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        textbox.text = self.content
        self.apply_font_style(textbox)
        textbox.text_frame.auto_size = True
        textbox.text_frame.word_wrap = True
        if self.location == 1:
            textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        elif self.location == 3:
            textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        else:
            textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        self.position_element(textbox)
        self.textbox = textbox

class CodeSnippet(Element):
    def __init__(self, content, style, bounding_box):
        super().__init__(content, style, bounding_box)

    def render(self, slide):
        left, top, width, height = self.bounding_box
        code_shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        code_tf = code_shape.text_frame
        if self.content != []:            
            code_tf.text = self.content[0]
            self.apply_font_style(code_shape)
            # enum_tf.paragraphs[0].paragraph_format.alignment = MSO_ANCHOR.JUSTIFY
            code_lines = self.content.split('\n')
            for i, pt_text in enumerate(code_lines):
                if i>0:
                    if isinstance(pt_text, str):
                        p = code_tf.add_paragraph()
                        run = p.add_run()
                        run.text = pt_text
                        self.apply_font_style_on_run(run)
                    else:
                        raise Exception("Invalid Code lines format")
                
        code_tf.auto_size = True
        code_tf.word_wrap = True
        MAX_LINES = 10
        if len(code_tf.paragraphs) > MAX_LINES:
            for paragraph in code_tf.paragraphs[MAX_LINES:]:
                paragraph.text = paragraph.text[:paragraph.text.rfind(' ')] + "..."
        self.position_element(code_shape)
        self.enum_tf = code_tf

class Graphic(Element):
    def __init__(self, content, style, bounding_box):
        super().__init__(content, style, bounding_box)
    
    def render(self, slide):
        left, top, width, height = self.bounding_box
        img = slide.shapes.add_picture(self.content, Inches(left), Inches(top), Inches(width), Inches(height))
        # slide.shapes._spTree.remove(img._element)
        # slide.shapes._spTree.insert(1, img._element)
        self.image = img 

class Template():
    def __init__(self, content):
        self.content = content
    def render(self, slide):
        img = slide.shapes.add_picture(self.content, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        slide.shapes._spTree.remove(img._element)
        slide.shapes._spTree.insert(2, img._element)  
        self.image = img

class PresentationGenerator:
    def __init__(self, json_payload, subject_name, slide_id, version):
        self.json_payload = json_payload
        self.subject_name = subject_name
        self.slide_id = slide_id
        self.version = version
        self.presentation = Presentation()
    
    def insert_title_slide(self, MUL_FAC=1):
        title_font = self.json_payload["slides"][0]["elements"]["title"][0]["style"]["font_name"]
        title_slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[0])
        title_shape = title_slide.shapes.title
        title_shape.text = self.json_payload['topic']
        title_shape.left = Inches(0.5)
        title_shape.top = Inches(2)
        title_shape.width = Inches(9*MUL_FAC)
        title_shape.height = Inches(1.25)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.size = Pt(48)
        title_shape.text_frame.paragraphs[0].font.name = title_font
        presenter = title_slide.shapes.add_textbox(Inches(3.5),Inches(3.25),Inches(3*MUL_FAC),Inches(0.75))
        presenter.text = self.json_payload["presenter"]
        presenter.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        presenter.text_frame.paragraphs[0].font.italic = True
        presenter.text_frame.paragraphs[0].font.size = Pt(28)
        presenter.text_frame.paragraphs[0].font.name = title_font
        date = title_slide.shapes.add_textbox(Inches(4),Inches(4),Inches(2*MUL_FAC),Inches(0.5))
        date.text = self.json_payload["date"]
        date.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        date.text_frame.paragraphs[0].font.name = title_font

    def generate_presentation(self):
        self.presentation.slide_width = Inches(13.333)
        self.presentation.slide_height = Inches(7.5)
        MUL_FAC = 1.33
        # self.insert_title_slide(MUL_FAC)
        for slide_info in self.json_payload['slides']:
            slide_layout = self.presentation.slide_layouts[1]
            slide = self.presentation.slides.add_slide(slide_layout)
            slide.background.fill.solid()
            if 'bg_color' in slide_info.keys():
                slide.background.fill.fore_color.rgb = RGBColor(slide_info['bg_color']['r'], slide_info['bg_color']['g'], slide_info['bg_color']['b'])
            else:
                slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)                
            
            if 'template' in slide_info.keys():
                element = Template(slide_info['template'])
                element.render(slide)

            for element_type, elements in slide_info['elements'].items():
                for element_info in elements:
                    if element_type == 'graphic':
                        if element_info['label'] == 'logo':
                            element = Graphic(element_info['value'], None, (element_info['xmin'], element_info['ymin'], element_info['width'], element_info['height']))
                    elif element_type == 'figures':
                        if 'caption' in element_info.keys():
                            element = Figure(element_info['path'], element_info['caption']['style'], (element_info['xmin'], element_info['ymin'], element_info['width'], element_info['height']), element_info['caption'])
                        else:
                            element = Figure(element_info['path'], None, (element_info['xmin'], element_info['ymin'], element_info['width'], element_info['height']), None)

                    elif element_type == 'equations':
                        if 'caption' in element_info.keys():
                            element = Equation(element_info['path'], element_info['caption']['style'], (element_info['xmin'], element_info['ymin'], element_info['width'], element_info['height']), element_info['caption'])
                        else:
                            element = Equation(element_info['path'], None, (element_info['xmin'], element_info['ymin'], element_info['width'], element_info['height']), None)
                    elif element_type == 'tables':
                        path = None
                        content = None
                        if 'path' in element_info.keys():
                            path = element_info['path']
                        else:
                            content = element_info['content']
                        if 'caption' in element_info.keys():
                            element = Table(path, element_info['caption']['style'], (element_info['xmin'], element_info['ymin'], element_info['width'], element_info['height']), element_info['caption'], content)
                        else:
                            element = Table(path, None, (element_info['xmin'], element_info['ymin'], element_info['width'], element_info['height']), None, content)
                    
                    elif element_type == 'code':
                        element = CodeSnippet(element_info['value'], element_info['style'], (element_info['xmin'], element_info['ymin'], element_info['width'], element_info['height']))
                    elif element_type == 'refs':
                        element = Reference(element_info['value'], element_info['style'], (element_info['xmin'], element_info['ymin'], element_info['width'], element_info['height']))
                    elif element_type == 'text':
                        if element_info['label'] == "enumeration":
                            if 'heading' in element_info.keys():
                                element = Enumeration(element_info['value'], element_info['style'], (element_info['xmin'], element_info['ymin'], element_info['width'], element_info['height']), element_info['heading'], None)
                            else:
                                element = Enumeration(element_info['value'], element_info['style'], (element_info['xmin'], element_info['ymin'], element_info['width'], element_info['height']), None, None)                             
                        else:
                            element = Description(element_info['value'], element_info['style'], (element_info['xmin'], element_info['ymin'], element_info['width'], element_info['height']), element_info["style"]["phrases"])
                    elif element_type == 'title':
                        element = Title(element_info['value'], element_info['style'], (element_info['xmin'], element_info['ymin'], element_info['width'], element_info['height']))
                    elif element_type == 'footer':
                        element = Footer(element_info['value'], element_info['style'], (element_info['xmin'], element_info['ymin'], element_info['width'], element_info['height']), element_info['location'])
                    else:
                        raise ValueError(f"Unsupported element type: {element_type}")

                    element.render(slide)

    
        # Clean-up empty elements
        for slide in self.presentation.slides:
            for shape in slide.shapes:
                # print(shape.shape_type)
                if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                    if not shape.text_frame.text:
                        sp = slide.shapes._spTree
                        sp.remove(shape._element)
        
        ppts_path = f"./ppts/"
        if os.path.isdir(ppts_path) == False:
            os.mkdir(ppts_path)

        if not os.path.exists(os.path.join(ppts_path, self.subject_name)):
            os.mkdir(os.path.join(ppts_path, self.subject_name))
        if not os.path.exists(os.path.join(ppts_path, self.subject_name, self.slide_id)):
            os.mkdir(os.path.join(ppts_path, self.subject_name, self.slide_id))

        self.presentation.save(os.path.join(ppts_path, self.subject_name, self.slide_id, f'{self.version}.pptx'))

def load_json_payload(file_path):
    with open(file_path, 'r') as file:
        return json.load(file)

def main():
    # Load the JSON payload
    print("Running PPT generator module...")
    buffer_folder_path = f"./code/buffer/full"
    entries = os.listdir(buffer_folder_path)
    # already generated presentations
    created_files = []
    for subject in os.listdir("dataset/json/"):
        for topic in os.listdir(f"dataset/json/{subject}"):
            created_files.append(topic)
    # Filter out directories
    directories = [entry for entry in entries if os.path.isdir(os.path.join(buffer_folder_path, entry))]
    json_file_paths = []
    for directory in directories:
        subject_dir = os.path.join(buffer_folder_path, directory)
        for root, dirs, files in os.walk(subject_dir):
            for file in files:
                if file.endswith('.json'):
                    json_file_paths.append(os.path.join(root, file))
    # print(directories)
    # for directory in directories:
    #     ppt_dirs = os.listdir(os.path.join(buffer_folder_path, directory))
    #     unique_ppt_dirs = [ppt_dir for ppt_dir in ppt_dirs if os.path.isdir(os.path.join(buffer_folder_path, directory, ppt_dir))]
        
    # json_file_paths = []
    # for ppt_dir in unique_ppt_dirs:
    #     json_files = [f for f in os.listdir(os.path.join(buffer_folder_path, directory)) if f.endswith('.json')]
    #     json_file_paths.append(os.path.join(buffer_folder_path, directory, json_files[-1]))
    

    base_topic_folder_path = f"./code/json"


    for i, json_file in enumerate(json_file_paths):
        if json_file.split("\\")[-2] in created_files:
            continue
        # Load JSON file from buffer
        subject_name = os.path.basename(os.path.dirname(os.path.dirname(json_file)))
        slide_id = os.path.basename(os.path.dirname(json_file))
        version , _ = os.path.splitext(os.path.basename(json_file))
        json_payload = load_json_payload(json_file)

        #Generate Presentation from JSON file and save it
        presentation_generator = PresentationGenerator(json_payload, subject_name, slide_id, version)
        presentation_generator.generate_presentation()
        # print(f"Presentation generated successfully for {slide_id}.")

    final_json_path = f"code/json/final"
    for i, json_file in enumerate(json_file_paths):
        if json_file.split("\\")[-2] in created_files:
            continue
        subject_name = os.path.basename(os.path.dirname(os.path.dirname(json_file)))
        slide_id = os.path.basename(os.path.dirname(json_file))
        version , _ = os.path.splitext(os.path.basename(json_file))
        if not os.path.exists(os.path.join(final_json_path, subject_name)):
            os.mkdir(os.path.join(final_json_path, subject_name))
        if not os.path.exists(os.path.join(final_json_path, subject_name, slide_id)):
            os.mkdir(os.path.join(final_json_path, subject_name, slide_id))
        final_json_file = os.path.join(final_json_path, subject_name, slide_id, f'{version}.json')
        
        if os.path.exists(json_file):
            os.rename(json_file, final_json_file)
    
    shutil.rmtree(f"code/buffer/temp")
    
    print("🟢 All presentations generated and files moved successfully.\n")

if __name__ == "__main__":
    main()