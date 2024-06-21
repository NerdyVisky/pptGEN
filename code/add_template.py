from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
import os
from random_generator import TEMPLATES
import random

light_colors = [
    (210, 218, 255),  # Alice Blue
    (225, 225, 200),  # Beige
    (255, 219, 183),  # Papaya Whip
    (255, 220, 175),  # Lemon Chiffon
    (210, 225, 210),  # Honeydew
    (235, 225, 218),  # Seashell
    (255, 228, 225),  # Misty Rose
    (245, 255, 250),  # Mint Cream
    (250, 240, 230),  # Linen
    (253, 245, 230)   # Old Lace
]

SOURCE_PATH = 'pptss'
changed = 0
total = 0
parent_dirs = os.listdir(SOURCE_PATH)
for parent_dir in parent_dirs:
    for item in os.listdir(os.path.join(SOURCE_PATH, parent_dir)):
        print(item)
        item_path = os.path.join(SOURCE_PATH, parent_dir, item)
        if os.path.isdir(item_path):
            vers = os.listdir(item_path)
            print(vers)
            for ver in vers:
                total += 1
                if random.random() > 0.8:
                    try:
                        presentation_path = os.path.join(SOURCE_PATH, parent_dir, item, ver)
                        print(presentation_path)
                        prs = Presentation(presentation_path)
                        image_path, isDark = TEMPLATES[random.randint(1, 8)]
                        for slide in prs.slides:
                            left = Inches(0)  
                            top = Inches(0)    
                            width = Inches(13.33)  
                            height = Inches(7.5) 
                            img = slide.shapes.add_picture(image_path, left, top, width, height)
                            slide.shapes._spTree.remove(img._element)
                            slide.shapes._spTree.insert(2, img._element)
                        prs.save(presentation_path)
                        changed += 1
                    except:
                        continue
                    continue

                if random.random() > 0.8:
                    try:
                        r, g, b = light_colors[random.randint(0, 9)]
                        # r, g, b = 255, 255, 255
                        presentation_path = os.path.join(SOURCE_PATH, parent_dir, item, ver)
                        print(presentation_path)
                        prs = Presentation(presentation_path)
                        image_path, isDark = TEMPLATES[random.randint(1, 8)]
                        for slide in prs.slides:
                            slide.background.fill.fore_color.rgb = RGBColor(r, g, b)
                        prs.save(presentation_path)
                        changed += 1
                    except:
                        continue
                    continue
                # output_path = 'path/to/your/modified_presentation.pptx'
                    
print(changed, total)
                    
                    
